from __future__ import annotations

from pathlib import Path
from typing import Union

import docx
from docx.oxml.table import CT_Tbl
from docx.oxml.text.paragraph import CT_P
from docx.table import Table
from docx.text.paragraph import Paragraph
from pptx import Presentation

PathLike = Union[str, Path]


def _iter_block_items(document: docx.document.Document):
    """Paragraphs and tables in document order (python-docx's recipe) — a มคอ's
    weekly-schedule/CLO tables sit between paragraphs and must stay positioned
    relative to their surrounding headings for the text to read sensibly."""
    for child in document.element.body.iterchildren():
        if isinstance(child, CT_P):
            yield Paragraph(child, document)
        elif isinstance(child, CT_Tbl):
            yield Table(child, document)


def _table_to_lines(table: Table) -> list:
    lines = []
    for row in table.rows:
        # horizontally merged cells repeat the same underlying cell several times
        # in row.cells — dedupe consecutive repeats instead of printing N copies.
        cells = [c.text.strip() for c in row.cells]
        deduped = []
        for cell in cells:
            if not deduped or deduped[-1] != cell:
                deduped.append(cell)
        if any(deduped):
            lines.append(" | ".join(deduped))
    return lines


def load_docx_text(path: PathLike) -> str:
    document = docx.Document(path)
    lines = []
    for block in _iter_block_items(document):
        if isinstance(block, Paragraph):
            text = block.text.strip()
            if text:
                lines.append(text)
        else:
            lines.extend(_table_to_lines(block))
    return "\n".join(lines)


def load_pptx_text(path: PathLike) -> str:
    prs = Presentation(path)
    slide_blocks = []
    for i, slide in enumerate(prs.slides, start=1):
        slide_lines = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    slide_lines.append(text)
        if slide_lines:  # skip blank/decorative divider slides
            slide_blocks.append(f"Slide {i}:\n" + "\n".join(slide_lines))
    return "\n\n".join(slide_blocks)


LOADERS = {".docx": load_docx_text, ".pptx": load_pptx_text}


def load_document_text(path: PathLike) -> str:
    suffix = Path(path).suffix.lower()
    loader = LOADERS.get(suffix)
    if loader is None:
        raise ValueError(f"Unsupported document type: {suffix or path}")
    return loader(path)
